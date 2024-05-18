import pandas as pd
from pptx import Presentation

game_class = int(input('Enter 1 for primary class.\nEnter 5 for middle class.\nEnter 8 for upper class.\n'))
class_column = f'Сложность для {game_class}-{game_class + 3} класса'

num_easy_questions = int(input('Enter number of easy questions: '))
num_medium_questions = int(input('Enter number of medium questions: '))
num_hard_questions = int(input('Enter number of hard questions: '))


excel_file_path = 'games.xlsx'
df = pd.read_excel(excel_file_path)


filtered_easy = df[(df[class_column] < 5) & (df[class_column] > 0)]
filtered_medium = df[(df[class_column] >= 5) & (df[class_column] < 8)]
filtered_hard = df[df[class_column] >= 8]

# проерка на достаточность вопросов
if len(filtered_easy) < num_easy_questions:
    print("Не хватает легких вопросов.")
    exit()
if len(filtered_medium) < num_medium_questions:
    print("Не хватает средних вопросов.")
    exit()
if len(filtered_hard) < num_hard_questions:
    print("Не хватает сложных вопросов.")
    exit()

# выбираем случайные вопросы
# sample выбирает случайные из filtered_easy количеством num_easy_questions
selected_easy = filtered_easy.sample(n=num_easy_questions)
selected_medium = filtered_medium.sample(n=num_medium_questions)
selected_hard = filtered_hard.sample(n=num_hard_questions)

# объединяем все выбранные вопросы
selected_questions = pd.concat([selected_easy, selected_medium, selected_hard])
# перемешиваем все вопросы и сбрасываем старые индексы
selected_questions = selected_questions.sample(frac=1).reset_index(drop=True)


prs = Presentation()
for index, row in selected_questions.iterrows():
    game_title = row['Название игры']
    question = row['Вопрос']
    answer = row['Ответ']

    # Создаем слайд с вопросом
    question_slide = prs.slides.add_slide(prs.slide_layouts[1])
    question_title = question_slide.shapes.title
    question_content = question_slide.placeholders[1]
    question_title.text = f"{game_title}"
    question_content.text = f"Вопрос: {question}"

    # Создаем слайд с ответом
    answer_slide = prs.slides.add_slide(prs.slide_layouts[1])
    answer_title = answer_slide.shapes.title
    answer_content = answer_slide.placeholders[1]
    answer_content.text = f"Ответ: {answer}"


prs.save('output_presentation.pptx')
